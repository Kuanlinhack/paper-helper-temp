import streamlit as st
import os
import tempfile
from PyPDF2 import PdfReader
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import yaml
import pymongo
import base64
from gridfs import GridFS
import hashlib

# streamlit run app_summary.py

# Configure Google Gemini API

# read api key from yaml file


# file_path = os.path.join(os.path.dirname(__file__), 'api_key.yaml')
# with open(file_path, 'r') as file:
#     api_key = yaml.safe_load(file)['Gemini_api_key']

api_key = "AIzaSyDBVn8ozBs1mX48xznLlRUIKy1c_LxqDK8"
genai.configure(api_key=api_key)
model = genai.GenerativeModel(model_name='gemini-1.5-flash')



def connect_to_mongodb(table_name="PaperHelper", collection_name="file", server_address = "mongodb://localhost:27017/"):
    """
    Connect to MongoDB and return the specified collection.
    
    Parameters:
        table_name (str): The name of the database to connect to.
        collection_name (str): The name of the collection to use.

    Returns:
        pymongo.collection.Collection: The MongoDB collection object.
    """
    try:
        # Connect to the MongoDB server
        client = pymongo.MongoClient(server_address)
        print("Connected to MongoDB successfully.")

        # Select the database (creates it if it doesn't exist)
        db = client[table_name]

        # Check if the database already exists
        if table_name in client.list_database_names():
            print(f"Database '{table_name}' already exists.")
        else:
            print(f"Database '{table_name}' does not exist. It will be created upon adding data.")

        # Select the collection (creates it if it doesn't exist)
        collection = db[collection_name]

        # Check if the collection already exists
        if collection_name in db.list_collection_names():
            print(f"Collection '{collection_name}' already exists.")
        else:
            print(f"Collection '{collection_name}' does not exist. It will be created upon adding data.")

        # Return the collection object for further use
        return collection
    except pymongo.errors.ConnectionError as e:
        print(f"Failed to connect to MongoDB: {e}")
        return None
    except Exception as e:
        print(f"An error occurred: {e}")
        return None

def read_pdf(pdf_path):
    """Read PDF and merge text content from all pages"""
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def generate_paper_title(pdf_text):
    """Generate paper title using Google Gemini API"""
    title_prompt = f"Extract or generate the most appropriate academic paper title from the following text. Only output the title:\n\n{pdf_text[:1000]}"
    try:
        title_response = model.generate_content(title_prompt)
        return title_response.text.strip() if title_response else "Untitled Paper"
    except Exception as e:
        st.error(f"Error generating title: {e}")
        return "Untitled Paper"

def generate_paper_summary(pdf_text):
    """Generate summary and key points using Google Gemini API"""
    summary_prompt = f"Please generate a comprehensive summary of the following PDF text, including main points and key insights. Provide the summary in clear, concise language:\n\n{pdf_text[:2000]}"
    try:
        summary_response = model.generate_content(summary_prompt)
        return summary_response.text.strip() if summary_response else "Unable to generate summary."
    except Exception as e:
        st.error(f"Error generating summary: {e}")
        return "An error occurred while generating the summary."



def generate_pdf_summary_and_details_for_ppt(pdf_text):
    """使用 Google Gemini API 生成 PDF 條列摘要與詳細重點整理"""
    summary_prompt = f"請為以下 PDF 文本生成條列式摘要，每點以「-」開頭，no more than 30 words， 約 ppt 半頁左右內容(用全英文）：\n\n{pdf_text[:2000]}"
    details_prompt = f"請為以下 PDF 文本生成更詳細的重點整理，不少於 500 字(用全英文）：\n\n{pdf_text[:2000]}"
    try:
        summary_response = model.generate_content(summary_prompt)
        summary = summary_response.text.strip() if summary_response else "未能生成摘要。"

        details_response = model.generate_content(details_prompt)
        details = details_response.text.strip() if details_response else "未能生成重點整理。"

        return summary, details
    except Exception as e:
        st.error(f"生成摘要與詳細時發生錯誤：{e}")
        return "未能生成摘要。", "未能生成詳細重點整理。"
 
def generate_image_description_for_page(image_path):
    """使用 Google Gemini API 上傳圖片並生成描述"""
    try:
        myfile = genai.upload_file(image_path)
        prompt = "Can you tell me about the photo? in concise way in Full English, 大概100字左右 ppt 半頁左右內容，用淺想易懂方式解釋"
        result = model.generate_content([myfile, "\n\n", prompt])
        return result.text.strip() if result else f"Unable to generate description for: {image_path}"
    except Exception as e:
        st.error(f"Error generating image description: {e}")
        return "Unable to generate image description"
    

def generate_image_description(image_path):
    """使用 Google Gemini API 上傳圖片並生成描述"""
    try:
        myfile = genai.upload_file(image_path)
        prompt = "Can you tell me about the photo? in concise way in Full English, 大概40字左右 ppt 半頁左右內容"
        result = model.generate_content([myfile, "\n\n", prompt])
        return result.text.strip() if result else f"Unable to generate description for: {image_path}"
    except Exception as e:
        st.error(f"Error generating image description: {e}")
        return "Unable to generate image description"




def add_image_and_description_to_slide(prs, image_path, description, max_chars_per_slide=400):
    """為每張圖片添加解說內容到單獨的幻燈片"""
    description_pages = [description[i:i + max_chars_per_slide] for i in range(0, len(description), max_chars_per_slide)]

    for page_num, page_content in enumerate(description_pages):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # 添加圖片
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(0.5), Inches(8), Inches(4.5))

        # 添加文字框
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(8), Inches(2))
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = text_frame.add_paragraph()
        p.text = page_content.strip()
        p.font.size = Pt(16)
        p.line_spacing = Pt(20)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

def add_text_slide(prs, title, content, bullet=False, font_size=20, max_chars_per_slide=800):
    """添加純文字幻燈片，支持分頁"""
    content_pages = [content[i:i + max_chars_per_slide] for i in range(0, len(content), max_chars_per_slide)]

    for page_num, page_content in enumerate(content_pages):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"{title}（第 {page_num + 1} 部分）"
        text_box = slide.placeholders[1]
        text_frame = text_box.text_frame
        text_frame.clear()

        if bullet:
            for line in page_content.split("\n"):
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(font_size)
        else:
            text_frame.text = page_content
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.alignment = PP_ALIGN.LEFT

def create_presentation(pdf_path, image_paths):
    """根據 PDF 和圖片生成 PPT，返回 PPT 檔案的路徑"""
    pdf_text = read_pdf(pdf_path)
    title = generate_paper_title(pdf_text)
    summary, details = generate_pdf_summary_and_details_for_ppt(pdf_text)

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "由 Paper Helper 生成的內容"

    add_text_slide(prs, "PDF 條列摘要", summary, bullet=True, font_size=20)
    add_text_slide(prs, "PDF 詳細整理", details, bullet=False, font_size=18)

    for image_path in image_paths:
        description = generate_image_description(image_path)
        add_image_and_description_to_slide(prs, image_path, description)

    output_path = os.path.join(tempfile.gettempdir(), f"{title}_presentation.pptx")
    prs.save(output_path)
    return output_path




def slides_for_ppt_app():

    # ======== Streamlit Web App ========

    st.markdown(
        """
        <style>
        .main-title {
            text-align: center;
            font-size: 36px;
            color: #4CAF50;
        }
        .sub-title {
            text-align: center;
            font-size: 18px;
            color: #555;
        }
        .upload-section {
            padding: 20px;
            background: #f9f9f9;
            border-radius: 10px;
            margin-bottom: 20px;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

    st.markdown('<h1 class="main-title">Paper Helper: Your Presentation Partner</h1>', unsafe_allow_html=True)
    st.markdown('<p class="sub-title">Upload your PDF and images to generate a stunning slides!</p>', unsafe_allow_html=True)

    # 上傳區域
    st.markdown('<div class="upload-section">', unsafe_allow_html=True)
    uploaded_pdf = st.file_uploader("選擇 PDF 文件", type=["pdf"])
    uploaded_images = st.file_uploader("選擇圖片（可多選）", type=["png", "jpg", "jpeg"], accept_multiple_files=True)
    st.markdown('</div>', unsafe_allow_html=True)

    if uploaded_pdf and uploaded_images and st.button("生成 PPT"):
        with st.spinner("正在生成簡報，請稍候..."):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                temp_pdf.write(uploaded_pdf.read())
                temp_pdf_path = temp_pdf.name

            image_paths = []
            for uploaded_image in uploaded_images:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_image:
                    temp_image.write(uploaded_image.read())
                    image_paths.append(temp_image.name)

            ppt_path = create_presentation(temp_pdf_path, image_paths)
            st.success("簡報已生成 🎉")

            with open(ppt_path, "rb") as ppt_file:
                st.download_button(
                    label="下載簡報",
                    data=ppt_file,
                    file_name=os.path.basename(ppt_path),
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

    # streamlit run  app.py


def image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def main():
    # Custom CSS for styling
    st.markdown("""
    <style>
    .main-title {
        text-align: center;
        font-size: 36px;
        color: #4CAF50;
    }
    .sub-title {
        text-align: center;
        font-size: 18px;
        color: #555;
    }
    .upload-section {
        padding: 20px;
        background: #f9f9f9;
        border-radius: 10px;
        margin-bottom: 20px;
    }
    .result-section {
        background: #f0f0f0;
        border-radius: 10px;
        padding: 20px;
        margin-top: 20px;
    }
    .paper-title {
        text-align: center;
        font-size: 24px;
        color: #333;
        margin-bottom: 15px;
    }
    </style>
    """, unsafe_allow_html=True)

    # App title and description
    st.markdown('<h1 class="main-title">Paper Helper</h1>', unsafe_allow_html=True)
    st.markdown('<p class="sub-title">Upload your PDF and images to get comprehensive insights</p>', unsafe_allow_html=True)

    # Upload section
    st.markdown('<div class="upload-section">', unsafe_allow_html=True)
    uploaded_pdf = st.file_uploader("Choose a PDF file", type=["pdf"])
    uploaded_images = st.file_uploader("Choose images (multiple allowed)", type=["png", "jpg", "jpeg"], accept_multiple_files=True)
    st.markdown('</div>', unsafe_allow_html=True)

    # Session state variables to retain results
    if "results" not in st.session_state:
        st.session_state["results"] = {}

    generate_insights_and_slides_button = st.button("Generate Insights and Slides")

    if uploaded_pdf and uploaded_images and generate_insights_and_slides_button:
        with st.spinner("Analyzing your document and images..."):
            # Save PDF to a temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                temp_pdf.write(uploaded_pdf.read())
                temp_pdf_path = temp_pdf.name

            # Read PDF content and generate title and summary
            pdf_text = read_pdf(temp_pdf_path)
            paper_title = generate_paper_title(pdf_text)
            summary = generate_paper_summary(pdf_text)

            # Store results in session state
            st.session_state["results"]["paper_title"] = paper_title
            st.session_state["results"]["summary"] = summary
            st.session_state["results"]["pdf_text"] = pdf_text

            # Generate image descriptions
            image_paths = []
            image_descriptions = []

            for uploaded_image in uploaded_images:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_image:
                    temp_image.write(uploaded_image.read())
                    image_path = temp_image.name
                    image_paths.append(image_path)
                    image_descriptions.append(generate_image_description_for_page(image_path))

            st.session_state["results"]["image_paths"] = image_paths
            st.session_state["results"]["image_descriptions"] = image_descriptions

            # Generate presentation
            ppt_path = create_presentation(temp_pdf_path, image_paths)
            st.session_state["results"]["ppt_path"] = ppt_path

            # connect to mongodb
            file_collection = connect_to_mongodb(table_name="PaperHelper", collection_name="file")
            # Initialize GridFS
            fs = GridFS(file_collection.database)

            # Generate a common ID based on the paper title
            paper_title = st.session_state["results"]["paper_title"]
            common_id = hashlib.md5(paper_title.encode()).hexdigest()

            # Combine all data to be saved in a single dictionary
            data_to_save = {
                "common_id": common_id,
                "paper_text": st.session_state["results"]["pdf_text"],
                "paper_title": paper_title,
                "image_paths": st.session_state["results"].get("image_paths", []),
                "image_descriptions": st.session_state["results"].get("image_descriptions", []),
                "images_base64": []
            }

            # Transform all images to base64 and add to the dictionary
            for image_path in data_to_save["image_paths"]:
                image_base64 = image_to_base64(image_path)
                data_to_save["images_base64"].append(image_base64)

            # Read the ppt file and save it to GridFS
            if "ppt_path" in st.session_state["results"]:
                with open(st.session_state["results"]["ppt_path"], "rb") as ppt_file:
                    ppt_file_id = fs.put(ppt_file, filename=os.path.basename(st.session_state["results"]["ppt_path"]))
                    data_to_save["ppt_file_id"] = ppt_file_id

            # Save the combined data to MongoDB
            file_collection.insert_one(data_to_save)

    # Display results if available in session state
    if "results" in st.session_state and "paper_title" in st.session_state["results"]:
        st.markdown('<div class="result-section">', unsafe_allow_html=True)
        st.markdown(f'<div class="paper-title" style="font-size: 28px;">Paper Title: {st.session_state["results"]["paper_title"]}</div>', unsafe_allow_html=True)
        st.subheader("Paper Summary")
        st.write(st.session_state["results"]["summary"])
        st.markdown('</div>', unsafe_allow_html=True)

        if "image_descriptions" in st.session_state["results"]:
            st.markdown('<div class="result-section">', unsafe_allow_html=True)
            st.subheader("Image Descriptions")
            for i, (image_path, description) in enumerate(zip(st.session_state["results"]["image_paths"], st.session_state["results"]["image_descriptions"])):
                st.markdown(f"### Image {i + 1}")
                st.image(image_path, use_container_width=True)
                st.write(description)
            st.markdown('</div>', unsafe_allow_html=True)

        st.success("Slides generated successfully 🎉")
        with open(st.session_state["results"]["ppt_path"], "rb") as ppt_file:
            st.download_button(
                label="Download Slides",
                data=ppt_file,
                file_name=os.path.basename(st.session_state["results"]["ppt_path"]),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )


if __name__ == "__main__":
    main()